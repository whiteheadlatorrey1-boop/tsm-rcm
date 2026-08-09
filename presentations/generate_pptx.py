#!/usr/bin/env python3
"""
Generates one leave-behind .pptx per vertical from the narrative data in
demo/presentations/data/*.json plus the REAL screenshots already captured
by the Playwright specs (tests/e2e/demo/screenshots/<folder>/<shot>.png).

Run this from inside the repo (Codespace), after screenshots exist:
    pip install python-pptx --break-system-packages   # if not already installed
    python3 demo/presentations/generate_pptx.py [vertical ...]

With no args, builds all 17. Pass one or more vertical keys (e.g. "schools")
to build just those. Output lands in demo/presentations/pptx/<vertical>-leave-behind.pptx

If a screenshot is missing, that slide gets a gray placeholder + warning text
instead of crashing the whole build — rerun the matching Playwright spec and
regenerate to fill it in.
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

BASE = os.path.dirname(os.path.abspath(__file__))
REPO_ROOT = os.path.abspath(os.path.join(BASE, "..", ".."))
DATA_DIR = os.path.join(BASE, "data")
SHOT_ROOT = os.path.join(REPO_ROOT, "tests", "e2e", "demo", "screenshots")
OUT_DIR = os.path.join(BASE, "pptx")
os.makedirs(OUT_DIR, exist_ok=True)

BG = RGBColor(0x07, 0x0D, 0x16)
PANEL = RGBColor(0x0E, 0x1A, 0x2B)
GOLD = RGBColor(0xE8, 0xB8, 0x4B)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
TEXT = RGBColor(0xE8, 0xEE, 0xF6)
MUTED = RGBColor(0x8A, 0xA5, 0xBF)
WARN = RGBColor(0xF8, 0x71, 0x71)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_bg(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_text(slide, left, top, width, height, text, size=18, color=TEXT,
             bold=False, align=PP_ALIGN.LEFT, font="Calibri", anchor=None, line_spacing=1.15):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font
    return box


def title_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(0.9), Inches(0.7), Inches(6), Inches(0.4),
              "TSM CONSULTZ", size=13, color=TEAL, bold=True)
    add_text(slide, Inches(0.9), Inches(2.3), Inches(11.5), Inches(1.8),
              data["title"], size=36, color=GOLD, bold=True)
    add_text(slide, Inches(0.9), Inches(4.25), Inches(11), Inches(0.8),
              data["tagline"], size=20, color=TEXT)
    pain_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                       Inches(0.9), Inches(5.3), Inches(11.3), Inches(1.4))
    pain_box.fill.solid(); pain_box.fill.fore_color.rgb = PANEL
    pain_box.line.color.rgb = RGBColor(0x16, 0x28, 0x3F)
    pain_box.line.width = Pt(1)
    tf = pain_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_top = Inches(0.2)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "THE PROBLEM"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = data["painPoint"]; r2.font.size = Pt(16)
    r2.font.color.rgb = TEXT


def step_slide(prs, data, step, idx, total, folder):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)

    # header
    add_text(slide, Inches(0.5), Inches(0.25), Inches(9), Inches(0.35),
              f"STEP {idx} OF {total}" + ("  \u2022  WOW MOMENT" if step.get("wow") else ""),
              size=11, color=(GOLD if step.get("wow") else TEAL), bold=True)
    add_text(slide, Inches(0.5), Inches(0.6), Inches(11.5), Inches(0.6),
              step.get("caption") or data["title"], size=24, color=TEXT, bold=True)

    # image (left ~60%)
    img_path = os.path.join(SHOT_ROOT, folder, step["shot"] + ".png")
    img_left, img_top = Inches(0.5), Inches(1.4)
    img_w, img_h = Inches(7.6), Inches(5.6)
    if os.path.exists(img_path):
        try:
            with Image.open(img_path) as im:
                iw, ih = im.size
            ratio = min(img_w / iw, img_h / ih)
            w, h = Emu(int(iw * ratio)), Emu(int(ih * ratio))
            left = Emu(int(img_left + (img_w - w) / 2))
            top = Emu(int(img_top + (img_h - h) / 2))
            slide.shapes.add_picture(img_path, left, top, width=w, height=h)
        except Exception as e:
            _placeholder(slide, img_left, img_top, img_w, img_h, f"Could not load image:\n{e}")
    else:
        _placeholder(slide, img_left, img_top, img_w, img_h,
                      f"Screenshot not yet captured.\nRun: npm run test:e2e -- {folder}\nExpected: {step['shot']}.png")

    if step.get("wow"):
        badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(0.7), Inches(1.55), Inches(1.9), Inches(0.4))
        badge.fill.solid(); badge.fill.fore_color.rgb = GOLD
        badge.line.fill.background()
        bt = badge.text_frame; bt.margin_left = Inches(0.05); bt.margin_top = Inches(0.02)
        bp = bt.paragraphs[0]; bp.alignment = PP_ALIGN.CENTER
        br = bp.add_run(); br.text = "\u2728 WOW MOMENT"; br.font.size = Pt(11)
        br.font.bold = True; br.font.color.rgb = RGBColor(0x1A, 0x14, 0x00)

    # talk track (right column)
    talk_left = Inches(8.3)
    talk_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, talk_left, Inches(1.4), Inches(4.5), Inches(5.6))
    talk_box.fill.solid(); talk_box.fill.fore_color.rgb = PANEL
    talk_box.line.fill.background()
    tf = talk_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.3); tf.margin_right = Inches(0.25); tf.margin_top = Inches(0.3)
    p = tf.paragraphs[0]
    r = p.add_run(); r.text = "SAY THIS:"; r.font.size = Pt(11); r.font.bold = True
    r.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    r2 = p2.add_run(); r2.text = step.get("talk", ""); r2.font.size = Pt(15)
    r2.font.color.rgb = TEXT
    p2.line_spacing = 1.3

    # speaker notes = full talk track too, for presenter view
    notes = slide.notes_slide.notes_text_frame
    notes.text = step.get("talk", "")


def _placeholder(slide, left, top, w, h, msg):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    box.fill.solid(); box.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x1A)
    box.line.color.rgb = WARN; box.line.width = Pt(1.5)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = msg; r.font.size = Pt(14); r.font.color.rgb = WARN


def closing_slide(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_bg(slide)
    add_text(slide, Inches(1.2), Inches(2.6), Inches(11), Inches(0.4),
              "THE CLOSE", size=13, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(3.1), Inches(11), Inches(2.2),
              data["cta"], size=30, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, Inches(1.2), Inches(5.6), Inches(11), Inches(0.5),
              "TSM CONSULTZ LLC \u2014 READY TO PILOT ON YOUR DATA", size=12, color=MUTED,
              align=PP_ALIGN.CENTER)


def build_one(key):
    data_path = os.path.join(DATA_DIR, f"{key}.json")
    if not os.path.exists(data_path):
        print(f"  ! no data file for '{key}', skipping"); return
    with open(data_path) as f:
        data = json.load(f)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    title_slide(prs, data)
    total = len(data["steps"])
    for i, step in enumerate(data["steps"], 1):
        if not step.get("talk"):
            continue  # skip demo-note-only pseudo-steps with no real content
        step_slide(prs, data, step, i, total, data["folder"])
    closing_slide(prs, data)

    out_path = os.path.join(OUT_DIR, f"{key}-leave-behind.pptx")
    prs.save(out_path)
    missing = sum(1 for s in data["steps"] if s.get("talk") and
                  not os.path.exists(os.path.join(SHOT_ROOT, data["folder"], s["shot"] + ".png")))
    flag = f"  ({missing} screenshot(s) missing \u2014 placeholders used)" if missing else ""
    print(f"  wrote {out_path}{flag}")


if __name__ == "__main__":
    keys = sys.argv[1:] or [f[:-5] for f in sorted(os.listdir(DATA_DIR)) if f.endswith(".json")]
    print(f"Building {len(keys)} deck(s)...")
    for k in keys:
        build_one(k)
    print("Done.")
